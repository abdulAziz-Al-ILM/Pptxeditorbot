import os
import logging
import asyncio
from telegram import Update
from telegram.ext import ApplicationBuilder, ContextTypes, CommandHandler, MessageHandler, filters
from pptx import Presentation
from pptx.util import Pt

# --- SOZLAMALAR ---
# Tokenni Railway Environment Variables dan oladi
BOT_TOKEN = os.getenv("BOT_TOKEN")

# Loglarni sozlash
logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)

# Foydalanuvchi holatini vaqtinchalik saqlash
user_data_store = {}

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Start komandasi."""
    await update.message.reply_text(
        "👋 Assalomu alaykum!\n\n"
        "Men PPTX faylning orqa fonini o'zgartirib beraman.\n"
        "1. Menga **.pptx** fayl yuboring.\n"
        "2. Keyin fon uchun **rasm** yuboring."
    )

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """PPTX faylni qabul qilish."""
    document = update.message.document
    file_name = document.file_name

    if not file_name.endswith('.pptx'):
        await update.message.reply_text("❌ Iltimos, faqat .pptx formatidagi fayl yuboring.")
        return

    user_id = update.effective_user.id
    file = await context.bot.get_file(document.file_id)
    
    # Faylni yuklab olish
    download_path = f"input_{user_id}.pptx"
    await file.download_to_drive(download_path)
    
    # Holatni saqlash
    user_data_store[user_id] = {'pptx_path': download_path}
    
    await update.message.reply_text(
        f"✅ {file_name} qabul qilindi.\n\n"
        "🖼 Endi yangi fon uchun **rasm** yuboring."
    )

async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Rasmni qabul qilish va yangi PPTX yaratish."""
    user_id = update.effective_user.id
    
    if user_id not in user_data_store or 'pptx_path' not in user_data_store[user_id]:
        await update.message.reply_text("⚠️ Iltimos, avval .pptx faylini yuboring.")
        return

    status_msg = await update.message.reply_text("⏳ Fayl tayyorlanmoqda...")

    # Rasmni yuklab olish
    photo_file = await context.bot.get_file(update.message.photo[-1].file_id)
    image_path = f"bg_{user_id}.jpg"
    await photo_file.download_to_drive(image_path)
    
    pptx_path = user_data_store[user_id]['pptx_path']
    output_pptx_path = f"yangi_slayd_{user_id}.pptx"

    try:
        # Sinxron funksiyani asinxron muhitda ishlatish (bloklamaslik uchun)
        loop = asyncio.get_running_loop()
        await loop.run_in_executor(None, create_new_presentation, pptx_path, image_path, output_pptx_path)
        
        # Faylni yuborish
        await update.message.reply_document(
            document=open(output_pptx_path, 'rb'),
            caption="✅ Marhamat, yangi slaydingiz!"
        )
        
    except Exception as e:
        logging.error(f"Xatolik: {e}")
        await update.message.reply_text("❌ Kechirasiz, faylni qayta ishlashda xatolik yuz berdi. Fayl buzilgan yoki juda murakkab bo'lishi mumkin.")
        
    finally:
        # Fayllarni tozalash
        for path in [pptx_path, image_path, output_pptx_path]:
            if os.path.exists(path):
                os.remove(path)
        
        if user_id in user_data_store:
            del user_data_store[user_id]
            
        try:
            await context.bot.delete_message(chat_id=update.effective_chat.id, message_id=status_msg.message_id)
        except:
            pass

def create_new_presentation(source_file, bg_image, output_file):
    """PPTX generator funksiyasi."""
    source_prs = Presentation(source_file)
    new_prs = Presentation()
    
    # O'lchamlarni nusxalash
    new_prs.slide_width = source_prs.slide_width
    new_prs.slide_height = source_prs.slide_height

    # Bo'sh slayd shabloni
    blank_layout = new_prs.slide_layouts[6]

    for slide in source_prs.slides:
        new_slide = new_prs.slides.add_slide(blank_layout)
        
        # 1. Fon rasmi
        new_slide.shapes.add_picture(bg_image, 0, 0, new_prs.slide_width, new_prs.slide_height)
        
        # 2. Matn qutilarini ko'chirish
        for shape in slide.shapes:
            if shape.has_text_frame:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                textbox = new_slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                
                # Agar matn bo'lsa
                if shape.text_frame.text.strip():
                    text_frame.clear()
                    
                    for paragraph in shape.text_frame.paragraphs:
                        new_p = text_frame.add_paragraph()
                        new_p.text = paragraph.text
                        
                        # Formatlashni saqlashga harakat qilamiz
                        if paragraph.font.size:
                            new_p.font.size = paragraph.font.size
                        else:
                            new_p.font.size = Pt(18) # Standart o'lcham
                        
                        new_p.font.bold = paragraph.font.bold
                        new_p.font.italic = paragraph.font.italic
                        # Rangni o'zgartirmaymiz (Original rang qoladi)

    new_prs.save(output_file)

if __name__ == '__main__':
    if not BOT_TOKEN:
        print("XATOLIK: BOT_TOKEN topilmadi. Railway Variables bo'limini tekshiring.")
    else:
        application = ApplicationBuilder().token(BOT_TOKEN).build()
        
        application.add_handler(CommandHandler("start", start))
        application.add_handler(MessageHandler(filters.Document.FileExtension("pptx"), handle_document))
        application.add_handler(MessageHandler(filters.PHOTO, handle_photo))
        
        print("Bot Railwayda ishga tushmoqda...")
        application.run_polling()
